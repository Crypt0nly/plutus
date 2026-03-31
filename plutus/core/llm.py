"""Model-agnostic LLM client powered by LiteLLM.

Supports Anthropic, OpenAI, local models (Ollama), and any OpenAI-compatible endpoint.
"""

from __future__ import annotations

import asyncio
import json
import logging
import os
import warnings
from collections.abc import AsyncIterator, Callable
from typing import Any

import litellm
from pydantic import BaseModel

from plutus.config import ModelConfig, SecretsStore

logger = logging.getLogger("plutus.llm")

# Suppress litellm's verbose logging
litellm.suppress_debug_info = True

# Drop unsupported params instead of raising errors.
# Some models (e.g. GPT-5) only accept temperature=1; rather than
# crashing, let litellm silently strip the unsupported kwarg.
litellm.drop_params = True

# Suppress Pydantic serialization warning for ServerToolUse.
# LiteLLM passes a plain dict where Pydantic expects a model instance when
# reporting server-side tool usage (web_search, web_fetch).  Harmless.
warnings.filterwarnings(
    "ignore",
    message=r".*Expected `ServerToolUse`.*",
    category=UserWarning,
)


class ToolDefinition(BaseModel):
    """Tool definition for function calling."""

    name: str
    description: str
    parameters: dict[str, Any]


class LLMOverloadedError(Exception):
    """Raised when all retries are exhausted due to provider overload (e.g. 529)."""

    def __init__(self, provider: str, original: Exception | None = None):
        self.provider = provider
        self.original = original
        super().__init__(
            f"{provider.capitalize()} is temporarily overloaded and all retries "
            f"have been exhausted. Original error: {original}"
        )


class ToolCall(BaseModel):
    """A tool call from the LLM response."""

    id: str
    name: str
    arguments: dict[str, Any]


class LLMMessage(BaseModel):
    role: str  # "system", "user", "assistant", "tool"
    content: str | None = None
    tool_calls: list[ToolCall] | None = None
    tool_call_id: str | None = None
    name: str | None = None


class LLMResponse(BaseModel):
    content: str | None = None
    tool_calls: list[ToolCall] = []
    finish_reason: str | None = None
    usage: dict[str, int] = {}


_ANTHROPIC_SERVER_TOOL_ID_PREFIX = "srvtoolu_"

# Anthropic server-side tool types (executed by Anthropic, not the client)
_ANTHROPIC_WEB_SEARCH_TOOL = {
    "type": "web_search_20250305",
    "name": "web_search",
    "max_uses": 5,
}
_ANTHROPIC_WEB_FETCH_TOOL = {
    "type": "web_fetch_20250910",
    "name": "web_fetch",
    "max_uses": 5,
}


# Supported image MIME types for multimodal messages
_IMAGE_TYPES = {"image/jpeg", "image/png", "image/gif", "image/webp"}

# MIME type for PDF documents (Anthropic-specific support)
_PDF_TYPE = "application/pdf"

# File types whose content can be extracted as text and sent inline to the LLM.
# Maps MIME type → human-readable label used in the text block header.
_TEXT_EXTRACTABLE: dict[str, str] = {
    # Plain text variants
    "text/plain": "Text file",
    "text/markdown": "Markdown document",
    "text/csv": "CSV data",
    "text/tsv": "TSV data",
    "text/tab-separated-values": "TSV data",
    "text/html": "HTML document",
    "text/css": "CSS stylesheet",
    "text/javascript": "JavaScript source",
    "text/x-python": "Python source",
    "text/x-script.python": "Python source",
    "text/x-java-source": "Java source",
    "text/x-c": "C source",
    "text/x-c++": "C++ source",
    "text/x-sh": "Shell script",
    "text/x-shellscript": "Shell script",
    "text/x-ruby": "Ruby source",
    "text/x-go": "Go source",
    "text/x-rust": "Rust source",
    "text/x-kotlin": "Kotlin source",
    "text/x-swift": "Swift source",
    "text/x-php": "PHP source",
    "text/x-sql": "SQL script",
    "text/x-yaml": "YAML document",
    "text/x-toml": "TOML document",
    "text/x-ini": "INI config",
    "text/x-log": "Log file",
    "text/x-diff": "Diff/patch",
    "text/x-patch": "Diff/patch",
    # Application text formats
    "application/json": "JSON data",
    "application/xml": "XML document",
    "application/x-yaml": "YAML document",
    "application/x-sh": "Shell script",
    "application/x-python-code": "Python source",
    "application/x-httpd-php": "PHP source",
    "application/javascript": "JavaScript source",
    "application/typescript": "TypeScript source",
    "application/toml": "TOML document",
    "application/sql": "SQL script",
    "application/graphql": "GraphQL schema",
    "application/ld+json": "JSON-LD data",
    "application/x-ndjson": "NDJSON data",
    "application/x-www-form-urlencoded": "Form data",
}

# Office document MIME types — require library-based extraction
_WORD_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # .docx
    "application/msword",  # .doc (legacy — best-effort)
}
_EXCEL_TYPES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # .xlsx
    "application/vnd.ms-excel",  # .xls (legacy)
    "application/vnd.oasis.opendocument.spreadsheet",  # .ods
}
_PPT_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
    "application/vnd.ms-powerpoint",  # .ppt (legacy)
    "application/vnd.oasis.opendocument.presentation",  # .odp
}
_CSV_TYPES = {
    "text/csv",
    "text/tsv",
    "text/tab-separated-values",
    "application/csv",
    "application/x-csv",
}
_ZIP_TYPES = {
    "application/zip",
    "application/x-zip-compressed",
    "application/x-zip",
    "application/octet-stream",  # generic binary — handled only if name ends .zip
}
_AUDIO_TYPES = {
    "audio/mpeg", "audio/mp3", "audio/wav", "audio/wave", "audio/ogg",
    "audio/flac", "audio/aac", "audio/webm", "audio/x-m4a", "audio/mp4",
}
_VIDEO_TYPES = {
    "video/mp4", "video/mpeg", "video/webm", "video/ogg", "video/quicktime",
    "video/x-msvideo", "video/x-matroska", "video/x-flv",
}

# Maximum characters of extracted text to include inline (prevents token bloat)
_MAX_TEXT_CHARS = 40_000


def _extract_docx(raw: bytes, name: str) -> str:
    """Extract text from a .docx file using python-docx."""
    try:
        import io
        import docx  # python-docx
        doc = docx.Document(io.BytesIO(raw))
        parts: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)
        text = "\n".join(parts)
        return text[:_MAX_TEXT_CHARS] if text else f"[{name}: document appears to be empty]"
    except ImportError:
        return f"[{name}: python-docx not installed — cannot extract Word document text]"
    except Exception as e:
        return f"[{name}: failed to extract Word document — {e}]"


def _extract_xlsx(raw: bytes, name: str) -> str:
    """Extract data from an Excel file using openpyxl."""
    try:
        import io
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"=== Sheet: {sheet_name} ===")
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    parts.append(",".join(cells))
                    row_count += 1
                if row_count >= 500:  # cap at 500 rows per sheet
                    parts.append("... [truncated — too many rows]")
                    break
        text = "\n".join(parts)
        return text[:_MAX_TEXT_CHARS] if text else f"[{name}: spreadsheet appears to be empty]"
    except ImportError:
        return f"[{name}: openpyxl not installed — cannot extract Excel data]"
    except Exception as e:
        return f"[{name}: failed to extract Excel data — {e}]"


def _extract_pptx(raw: bytes, name: str) -> str:
    """Extract slide text from a .pptx file using python-pptx."""
    try:
        import io
        import pptx  # python-pptx
        prs = pptx.Presentation(io.BytesIO(raw))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_texts.append(t)
            if slide_texts:
                parts.append(f"--- Slide {i} ---")
                parts.extend(slide_texts)
        text = "\n".join(parts)
        return text[:_MAX_TEXT_CHARS] if text else f"[{name}: presentation appears to have no text]"
    except ImportError:
        return f"[{name}: python-pptx not installed — cannot extract PowerPoint text]"
    except Exception as e:
        return f"[{name}: failed to extract PowerPoint text — {e}]"


def _extract_zip_listing(raw: bytes, name: str) -> str:
    """List contents of a ZIP archive and extract small text files."""
    try:
        import io
        import zipfile
        zf = zipfile.ZipFile(io.BytesIO(raw))
        entries = zf.infolist()
        parts = [f"ZIP archive: {name} ({len(entries)} entries)"]
        extracted_texts: list[str] = []
        total_extracted = 0
        for entry in entries[:200]:  # cap listing at 200 entries
            size_kb = entry.file_size / 1024
            parts.append(f"  {entry.filename} ({size_kb:.1f} KB)")
            # Auto-extract small text files
            if (
                not entry.is_dir()
                and entry.file_size < 50_000
                and any(entry.filename.endswith(ext) for ext in (
                    ".txt", ".md", ".py", ".js", ".ts", ".json", ".yaml",
                    ".yml", ".toml", ".csv", ".xml", ".html", ".css",
                    ".sh", ".bash", ".env", ".gitignore", ".dockerfile",
                ))
            ):
                try:
                    content = zf.read(entry.filename).decode("utf-8", errors="replace")
                    extracted_texts.append(f"\n--- {entry.filename} ---\n{content[:2000]}")
                    total_extracted += len(content)
                    if total_extracted > 20_000:
                        break
                except Exception:
                    pass
        if len(entries) > 200:
            parts.append(f"  ... and {len(entries) - 200} more entries")
        if extracted_texts:
            parts.append("\nExtracted text files:")
            parts.extend(extracted_texts)
        return "\n".join(parts)[:_MAX_TEXT_CHARS]
    except Exception as e:
        return f"[{name}: failed to read ZIP archive — {e}]"

# OpenAI models that support native computer use via the Responses API
_OPENAI_COMPUTER_USE_MODELS = {"gpt-5.4", "computer-use-preview"}

# Special tool call name used when the model natively returns a computer_call.
# The agent loop checks for this name to execute via ComputerUseExecutor
# instead of the normal tool registry.
NATIVE_COMPUTER_USE_TOOL = "__computer_use__"


class LLMClient:
    """Unified LLM client with tool-calling support."""

    def __init__(self, config: ModelConfig, secrets: SecretsStore | None = None):
        self._config = config
        self._secrets = secrets or SecretsStore()
        self._model = self._resolve_model()
        self._key_available = self._ensure_api_key()

    @property
    def key_configured(self) -> bool:
        """Whether an API key is available for the current provider."""
        return self._key_available

    def _resolve_model(self) -> str:
        """Map provider/model to litellm model string."""
        provider = self._config.provider
        model = self._config.model

        if provider == "anthropic":
            return f"anthropic/{model}" if not model.startswith("anthropic/") else model
        if provider == "openai":
            return f"openai/{model}" if not model.startswith("openai/") else model
        if provider == "ollama":
            return f"ollama/{model}" if not model.startswith("ollama/") else model
        # Custom / OpenAI-compatible
        return model

    def _ensure_api_key(self) -> bool:
        """Resolve the API key from env var or secrets store.

        Returns True if a key is available, False otherwise.
        Does NOT crash — the server can start without a key and prompt the user.
        """
        if self._config.provider in ("ollama", "local"):
            return True

        # Try secrets store (checks env var first, then file)
        key = self._secrets.get_key(self._config.provider)
        if key:
            # Always overwrite os.environ so a stale or wrong key from a previous
            # provider never blocks the real key.  The startup inject_all() guard
            # ("don't overwrite existing env vars") is intentional for startup only;
            # here we are explicitly reloading so we must force the update.
            env_var = self._config.api_key_env
            os.environ[env_var] = key
            return True

        logger.warning(
            f"No API key found for {self._config.provider}. "
            f"Set {self._config.api_key_env} or use the web UI to configure."
        )
        return False

    def reload_model(self, config: ModelConfig) -> None:
        """Hot-reload model configuration (called after user changes model via UI)."""
        self._config = config
        self._model = self._resolve_model()
        self._key_available = self._ensure_api_key()
        logger.info(f"Model reloaded: {self._model} (provider={config.provider})")

    def reload_key(self) -> bool:
        """Re-check for API key availability (called after user sets a key via UI)."""
        self._key_available = self._ensure_api_key()
        return self._key_available

    @property
    def is_anthropic(self) -> bool:
        """Whether the current model is an Anthropic model."""
        return self._config.provider == "anthropic" or self._model.startswith("anthropic/")

    @property
    def is_openai(self) -> bool:
        """Whether the current model is an OpenAI model."""
        return self._config.provider == "openai" or self._model.startswith("openai/")

    @property
    def supports_native_computer_use(self) -> bool:
        """Whether the current model supports OpenAI's native computer use tool.

        When True, the agent loop should handle ``__computer_use__`` tool calls
        by executing actions via ComputerUseExecutor instead of delegating to
        the ``openai_computer`` wrapper tool.
        """
        if not self.is_openai:
            return False
        return self._config.model.lower() in _OPENAI_COMPUTER_USE_MODELS

    def _build_kwargs(self, **overrides: Any) -> dict[str, Any]:
        kwargs: dict[str, Any] = {
            "model": self._model,
            "temperature": self._config.temperature,
            "max_tokens": self._config.max_tokens,
        }
        if self._config.base_url:
            kwargs["api_base"] = self._config.base_url
        # Explicitly pass the API key to litellm so it never falls back to a
        # stale or wrong key that may be sitting in os.environ from a previous
        # provider.  Only inject for providers that need a key.
        if self._config.provider not in ("ollama", "local"):
            api_key = self._secrets.get_key(self._config.provider)
            if api_key:
                kwargs["api_key"] = api_key
        kwargs.update(overrides)
        return kwargs

    def _expand_attachments(self, messages: list[dict[str, Any]]) -> list[dict[str, Any]]:
        """Convert messages with 'attachments' into multimodal content blocks.

        Handles both Anthropic and OpenAI formats:
        - Images: Use OpenAI-style image_url blocks (LiteLLM translates for Anthropic)
        - PDFs (Anthropic only): Use document source blocks
        - Word/Excel/PowerPoint: Extract text via python-docx / openpyxl / python-pptx
        - CSV, JSON, plain text, code files: Decode UTF-8 and include inline
        - ZIP archives: List contents and extract small text files
        - Audio/Video: Include metadata note (cannot embed in LLM context)
        """
        import base64 as _b64

        result = []
        for msg in messages:
            attachments = msg.pop("attachments", None)
            if not attachments or msg.get("role") != "user":
                result.append(msg)
                continue

            # Build multimodal content array
            content_blocks: list[dict[str, Any]] = []

            # Add text block first (if any)
            text = msg.get("content", "")
            if text:
                content_blocks.append({"type": "text", "text": text})

            for att in attachments:
                mime = att.get("type", "") or ""
                data = att.get("data", "")  # base64-encoded
                name = att.get("name", "file")
                # Normalise MIME — browsers sometimes send empty string for
                # unknown types; fall back to extension-based detection.
                if not mime or mime == "application/octet-stream":
                    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
                    mime = {
                        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        "doc":  "application/msword",
                        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        "xls":  "application/vnd.ms-excel",
                        "ods":  "application/vnd.oasis.opendocument.spreadsheet",
                        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        "ppt":  "application/vnd.ms-powerpoint",
                        "odp":  "application/vnd.oasis.opendocument.presentation",
                        "pdf":  "application/pdf",
                        "zip":  "application/zip",
                        "csv":  "text/csv",
                        "tsv":  "text/tsv",
                        "json": "application/json",
                        "xml":  "application/xml",
                        "yaml": "text/x-yaml",
                        "yml":  "text/x-yaml",
                        "toml": "application/toml",
                        "md":   "text/markdown",
                        "txt":  "text/plain",
                        "py":   "text/x-python",
                        "js":   "text/javascript",
                        "ts":   "application/typescript",
                        "jsx":  "text/javascript",
                        "tsx":  "application/typescript",
                        "html": "text/html",
                        "css":  "text/css",
                        "sh":   "text/x-sh",
                        "bash": "text/x-sh",
                        "sql":  "application/sql",
                        "rs":   "text/x-rust",
                        "go":   "text/x-go",
                        "java": "text/x-java-source",
                        "kt":   "text/x-kotlin",
                        "swift":"text/x-swift",
                        "rb":   "text/x-ruby",
                        "php":  "text/x-php",
                        "c":    "text/x-c",
                        "cpp":  "text/x-c++",
                        "h":    "text/x-c",
                        "log":  "text/x-log",
                        "diff": "text/x-diff",
                        "patch":"text/x-patch",
                        "mp3":  "audio/mpeg",
                        "wav":  "audio/wav",
                        "ogg":  "audio/ogg",
                        "flac": "audio/flac",
                        "aac":  "audio/aac",
                        "m4a":  "audio/x-m4a",
                        "mp4":  "video/mp4",
                        "mov":  "video/quicktime",
                        "avi":  "video/x-msvideo",
                        "mkv":  "video/x-matroska",
                        "webm": "video/webm",
                    }.get(ext, "application/octet-stream")

                # ── Images ──────────────────────────────────────────────────
                if mime in _IMAGE_TYPES:
                    content_blocks.append({
                        "type": "image_url",
                        "image_url": {"url": f"data:{mime};base64,{data}"},
                    })

                # ── PDFs ────────────────────────────────────────────────────
                elif mime == _PDF_TYPE and self.is_anthropic:
                    content_blocks.append({
                        "type": "document",
                        "source": {"type": "base64", "media_type": _PDF_TYPE, "data": data},
                    })
                elif mime == _PDF_TYPE:
                    content_blocks.append({"type": "text", "text": (
                        f"[Attached PDF: {name} — PDF rendering is only supported with "
                        f"Anthropic models. Switch to Claude to read this file.]"
                    )})

                # ── Word documents ───────────────────────────────────────────
                elif mime in _WORD_TYPES:
                    raw = _b64.b64decode(data)
                    extracted = _extract_docx(raw, name)
                    content_blocks.append({"type": "text", "text": (
                        f"=== Word document: {name} ===\n{extracted}"
                    )})

                # ── Excel spreadsheets ──────────────────────────────────────
                elif mime in _EXCEL_TYPES:
                    raw = _b64.b64decode(data)
                    extracted = _extract_xlsx(raw, name)
                    content_blocks.append({"type": "text", "text": (
                        f"=== Excel spreadsheet: {name} ===\n{extracted}"
                    )})

                # ── PowerPoint presentations ─────────────────────────────────
                elif mime in _PPT_TYPES:
                    raw = _b64.b64decode(data)
                    extracted = _extract_pptx(raw, name)
                    content_blocks.append({"type": "text", "text": (
                        f"=== PowerPoint presentation: {name} ===\n{extracted}"
                    )})

                # ── ZIP archives ─────────────────────────────────────────────
                elif mime in _ZIP_TYPES or name.lower().endswith(".zip"):
                    raw = _b64.b64decode(data)
                    extracted = _extract_zip_listing(raw, name)
                    content_blocks.append({"type": "text", "text": extracted})

                # ── Audio / Video — metadata note only ──────────────────────
                elif mime in _AUDIO_TYPES:
                    size_kb = len(data) * 3 // 4 // 1024  # approx decoded size
                    content_blocks.append({"type": "text", "text": (
                        f"[Attached audio file: {name} (~{size_kb} KB, {mime}). "
                        f"The file was uploaded. To transcribe or analyse it, "
                        f"save it to disk and use a speech-to-text tool.]"
                    )})
                elif mime in _VIDEO_TYPES:
                    size_kb = len(data) * 3 // 4 // 1024
                    content_blocks.append({"type": "text", "text": (
                        f"[Attached video file: {name} (~{size_kb} KB, {mime}). "
                        f"The file was uploaded. To process it, save it to disk "
                        f"and use a video analysis tool.]"
                    )})

                # ── Plain text, code, CSV, JSON, XML, YAML, etc. ─────────────
                elif mime in _TEXT_EXTRACTABLE or mime.startswith("text/"):
                    try:
                        raw = _b64.b64decode(data)
                        text_content = raw.decode("utf-8", errors="replace")
                        label = _TEXT_EXTRACTABLE.get(mime, "Text file")
                        if len(text_content) > _MAX_TEXT_CHARS:
                            text_content = text_content[:_MAX_TEXT_CHARS] + "\n... [truncated]"
                        content_blocks.append({"type": "text", "text": (
                            f"=== {label}: {name} ===\n{text_content}"
                        )})
                    except Exception as exc:
                        content_blocks.append({"type": "text", "text": (
                            f"[{name}: could not decode as text — {exc}]"
                        )})

                # ── Unknown / binary ─────────────────────────────────────────
                else:
                    size_kb = len(data) * 3 // 4 // 1024
                    content_blocks.append({"type": "text", "text": (
                        f"[Attached file: {name} ({mime}, ~{size_kb} KB) — "
                        f"binary format not directly readable. "
                        f"Save it to disk to process it with an appropriate tool.]"
                    )})

            # Replace content with multimodal blocks
            new_msg = {k: v for k, v in msg.items() if k != "content"}
            new_msg["content"] = content_blocks
            result.append(new_msg)

        return result

    @staticmethod
    def _sanitize_messages(messages: list[dict[str, Any]]) -> list[dict[str, Any]]:
        """Ensure every tool_use has a matching tool_result immediately after.

        Anthropic requires strict pairing: each assistant message with tool_calls
        must be followed by tool result messages for every tool_call_id.
        If any are missing, we inject synthetic tool_result messages to prevent
        API errors.
        """
        sanitized: list[dict[str, Any]] = []
        for i, msg in enumerate(messages):
            sanitized.append(msg)

            # Check if this is an assistant message with tool_calls
            tool_calls = msg.get("tool_calls", [])
            if msg.get("role") == "assistant" and tool_calls:
                # Collect the tool_call IDs that need results.
                # Skip server-side tool IDs (srvtoolu_*) — those are
                # handled by Anthropic's servers and don't need client results.
                expected_ids = set()
                for tc in tool_calls:
                    tc_id = tc.get("id")
                    if tc_id and not tc_id.startswith(_ANTHROPIC_SERVER_TOOL_ID_PREFIX):
                        expected_ids.add(tc_id)

                # Look ahead for tool_result messages that follow
                found_ids = set()
                for j in range(i + 1, len(messages)):
                    next_msg = messages[j]
                    if next_msg.get("role") == "tool":
                        tcid = next_msg.get("tool_call_id")
                        if tcid in expected_ids:
                            found_ids.add(tcid)
                    else:
                        break  # Stop at first non-tool message

                # Inject synthetic results for any missing tool_call_ids
                missing = expected_ids - found_ids
                for tc_id in missing:
                    logger.warning(f"Injecting synthetic tool_result for orphaned tool_use {tc_id}")
                    # We need to insert AFTER the current message but BEFORE
                    # the next non-tool message. Since we're building a new list,
                    # we can't insert into `sanitized` mid-iteration easily.
                    # Instead, we'll do a second pass below.

        # Second pass: ensure pairing
        result: list[dict[str, Any]] = []
        i = 0
        while i < len(sanitized):
            msg = sanitized[i]
            result.append(msg)

            tool_calls = msg.get("tool_calls", [])
            if msg.get("role") == "assistant" and tool_calls:
                expected_ids = set()
                for tc in tool_calls:
                    tc_id = tc.get("id")
                    if tc_id and not tc_id.startswith(_ANTHROPIC_SERVER_TOOL_ID_PREFIX):
                        expected_ids.add(tc_id)

                # Consume following tool messages
                j = i + 1
                while j < len(sanitized) and sanitized[j].get("role") == "tool":
                    tool_msg = sanitized[j]
                    # Anthropic requires tool messages to always have non-empty content
                    if not tool_msg.get("content"):
                        tool_msg["content"] = "(no output)"
                    result.append(tool_msg)
                    tcid = tool_msg.get("tool_call_id")
                    if tcid in expected_ids:
                        expected_ids.discard(tcid)
                    j += 1

                # Inject synthetic results for any still-missing IDs
                for tc_id in expected_ids:
                    logger.warning(f"Sanitizer: injecting tool_result for orphaned tool_use {tc_id}")
                    result.append({
                        "role": "tool",
                        "tool_call_id": tc_id,
                        "content": "[No result available — tool execution was interrupted]",
                    })

                i = j
                continue

            i += 1

        # Final sweep: ensure every tool message has non-empty content and every
        # assistant message has a content field.  This catches edge cases where
        # tool messages end up in unexpected positions (e.g. after context-window
        # truncation splits an assistant+tool pair).
        for msg in result:
            if msg.get("role") == "tool":
                if not msg.get("content"):
                    msg["content"] = "(no output)"
            elif msg.get("role") == "assistant":
                # Anthropic requires a content field on assistant messages even
                # when the message only contains tool_use blocks.
                if "content" not in msg:
                    msg["content"] = ""

        return result

    def _build_tools_list(
        self, tools: list[ToolDefinition] | None
    ) -> list[dict[str, Any]] | None:
        """Build the tools list, injecting Anthropic server-side tools when applicable.

        OpenAI models use the flat format (Responses API compatible):
            {"type": "function", "name": ..., "description": ..., "parameters": ...}
        Other providers use the nested Chat Completions format:
            {"type": "function", "function": {"name": ..., ...}}
        """
        if not tools:
            # Even with no function tools, inject server-side tools for Anthropic
            if self.is_anthropic and self._config.web_search:
                return [_ANTHROPIC_WEB_SEARCH_TOOL, _ANTHROPIC_WEB_FETCH_TOOL]
            return None

        if self.is_openai:
            # OpenAI Responses API flat format — also accepted by Chat
            # Completions via LiteLLM translation
            tool_list: list[dict[str, Any]] = [
                {
                    "type": "function",
                    "name": t.name,
                    "description": t.description,
                    "parameters": t.parameters,
                }
                for t in tools
            ]
        else:
            # Chat Completions nested format (Anthropic, Ollama, etc.)
            tool_list = [
                {
                    "type": "function",
                    "function": {
                        "name": t.name,
                        "description": t.description,
                        "parameters": t.parameters,
                    },
                }
                for t in tools
            ]

        # Inject Anthropic server-side web search/fetch tools.
        # These are executed by Anthropic's servers — the client never
        # handles the tool call; results come back inside the response.
        if self.is_anthropic and self._config.web_search:
            tool_list.append(_ANTHROPIC_WEB_SEARCH_TOOL)
            tool_list.append(_ANTHROPIC_WEB_FETCH_TOOL)

        return tool_list

    async def complete(
        self,
        messages: list[dict[str, Any]],
        tools: list[ToolDefinition] | None = None,
        on_retry: Callable[[int, int, str], Any] | None = None,
        **kwargs: Any,
    ) -> LLMResponse:
        """Send a completion request to the LLM.

        When the model supports native computer use (e.g. GPT-5.4), this
        routes through the OpenAI Responses API instead of LiteLLM so that
        the ``{"type": "computer"}`` tool is available natively.

        Args:
            on_retry: Optional callback ``(attempt, wait_seconds, message)``
                called before each retry sleep so the caller can surface
                progress to the user.
        """
        if self.supports_native_computer_use:
            return await self._complete_openai_native(messages, tools, **kwargs)

        call_kwargs = self._build_kwargs(**kwargs)
        expanded = self._expand_attachments(messages)
        call_kwargs["messages"] = self._sanitize_messages(expanded)

        built_tools = self._build_tools_list(tools)
        if built_tools:
            call_kwargs["tools"] = built_tools

        # Retry up to 4 times on transient server errors (Anthropic 500, 529, etc.)
        # Longer waits give Anthropic time to recover from brief outages.
        _retryable = (
            litellm.InternalServerError,
            litellm.ServiceUnavailableError,
            litellm.RateLimitError,
        )
        _retry_waits = [2, 5, 10, 20]  # seconds between attempts
        last_exc: Exception | None = None
        for attempt in range(len(_retry_waits) + 1):
            try:
                response = await litellm.acompletion(**call_kwargs)
                return self._parse_response(response)
            except _retryable as exc:
                last_exc = exc
                if attempt >= len(_retry_waits):
                    break  # exhausted all retries
                wait = _retry_waits[attempt]
                provider = self._config.provider.capitalize()
                msg = (
                    f"{provider} is temporarily overloaded. "
                    f"Retrying in {wait}s (attempt {attempt + 1}/{len(_retry_waits) + 1})..."
                )
                logger.warning(msg)
                if on_retry:
                    try:
                        result = on_retry(attempt, wait, msg)
                        if asyncio.iscoroutine(result):
                            await result
                    except Exception:
                        pass  # never let callback errors break the retry loop
                await asyncio.sleep(wait)

        # Wrap the final exception with provider info for fallback handling
        raise LLMOverloadedError(
            provider=self._config.provider,
            original=last_exc,
        )

    async def stream(
        self,
        messages: list[dict[str, Any]],
        tools: list[ToolDefinition] | None = None,
        **kwargs: Any,
    ) -> AsyncIterator[str]:
        """Stream a completion response, yielding text chunks."""
        call_kwargs = self._build_kwargs(stream=True, **kwargs)
        call_kwargs["messages"] = self._sanitize_messages(messages)

        built_tools = self._build_tools_list(tools)
        if built_tools:
            call_kwargs["tools"] = built_tools

        response = await litellm.acompletion(**call_kwargs)
        async for chunk in response:
            delta = chunk.choices[0].delta
            if delta.content:
                yield delta.content

    async def _complete_openai_native(
        self,
        messages: list[dict[str, Any]],
        tools: list[ToolDefinition] | None = None,
        **kwargs: Any,
    ) -> LLMResponse:
        """Send a completion via the OpenAI Responses API with native computer use.

        This bypasses LiteLLM and uses the OpenAI SDK directly so we can
        include ``{"type": "computer"}`` alongside regular function tools.
        """
        from openai import AsyncOpenAI

        api_key = os.environ.get(self._config.api_key_env)
        client = AsyncOpenAI(api_key=api_key)

        # Convert our message format to Responses API input items
        instructions = None
        input_items: list[dict[str, Any]] = []

        # Track which tool_call IDs were computer calls so we can format
        # their results as computer_call_output instead of function_call_output.
        computer_call_ids: set[str] = set()

        for msg in messages:
            role = msg.get("role")

            if role == "system":
                instructions = msg.get("content", "")

            elif role == "user":
                content = msg.get("content", "")
                if isinstance(content, str):
                    input_items.append({"role": "user", "content": content})
                else:
                    # Multimodal content blocks — pass through
                    input_items.append({"role": "user", "content": content})

            elif role == "assistant":
                content = msg.get("content")
                tool_calls = msg.get("tool_calls") or []

                # Emit text content as an assistant message
                if content:
                    input_items.append({
                        "role": "assistant",
                        "content": [{"type": "output_text", "text": content}],
                    })

                # Emit tool calls in their native format.
                # Messages from history use Chat Completions format where
                # name/arguments are nested under "function"; extract them.
                for tc in tool_calls:
                    tc_id = tc.get("id", "")
                    tc_name = tc.get("name", "")
                    tc_args = tc.get("arguments", {})
                    if not tc_name and "function" in tc:
                        func = tc["function"]
                        tc_name = func.get("name", "")
                        tc_args = func.get("arguments", tc_args)
                        if isinstance(tc_args, str):
                            try:
                                tc_args = json.loads(tc_args)
                            except (json.JSONDecodeError, ValueError):
                                pass

                    if tc_name == NATIVE_COMPUTER_USE_TOOL:
                        # Reconstruct the computer_call item
                        computer_call_ids.add(tc_id)
                        input_items.append({
                            "type": "computer_call",
                            "call_id": tc_id,
                            "actions": tc_args.get("actions", []),
                        })
                    else:
                        # Regular function call
                        args_str = json.dumps(tc_args) if isinstance(tc_args, dict) else tc_args
                        input_items.append({
                            "type": "function_call",
                            "call_id": tc_id,
                            "name": tc_name,
                            "arguments": args_str,
                        })

            elif role == "tool":
                tc_id = msg.get("tool_call_id", "")
                content = msg.get("content", "")

                if tc_id in computer_call_ids:
                    # Extract screenshot URL from the JSON result
                    screenshot_url = ""
                    try:
                        data = json.loads(content)
                        screenshot_url = data.get("screenshot_url", "")
                    except (json.JSONDecodeError, ValueError):
                        screenshot_url = content

                    input_items.append({
                        "type": "computer_call_output",
                        "call_id": tc_id,
                        "output": {
                            "type": "computer_screenshot",
                            "image_url": screenshot_url,
                        },
                    })
                else:
                    input_items.append({
                        "type": "function_call_output",
                        "call_id": tc_id,
                        "output": content,
                    })

        # Build tools: native computer + function tools
        # Responses API uses flat format (no nested "function" key)
        api_tools: list[dict[str, Any]] = [{"type": "computer"}]
        if tools:
            for t in tools:
                api_tools.append({
                    "type": "function",
                    "name": t.name,
                    "description": t.description,
                    "parameters": t.parameters,
                })

        create_kwargs: dict[str, Any] = {
            "model": self._config.model,
            "tools": api_tools,
            "input": input_items,
            "truncation": "auto",
        }
        if instructions:
            create_kwargs["instructions"] = instructions

        response = await client.responses.create(**create_kwargs)
        return self._parse_openai_native_response(response)

    def _parse_openai_native_response(self, response: Any) -> LLMResponse:
        """Parse an OpenAI Responses API response into our LLMResponse format."""
        content_parts: list[str] = []
        tool_calls: list[ToolCall] = []

        for item in response.output:
            item_type = getattr(item, "type", None)

            if item_type == "message":
                for block in getattr(item, "content", []):
                    if getattr(block, "type", None) == "output_text":
                        content_parts.append(block.text)

            elif item_type == "computer_call":
                # Convert to a ToolCall with our special name so the agent
                # loop knows to execute via ComputerUseExecutor.
                actions_raw = getattr(item, "actions", []) or []
                actions = []
                for a in actions_raw:
                    if hasattr(a, "model_dump"):
                        actions.append(a.model_dump())
                    elif hasattr(a, "__dict__"):
                        actions.append(
                            {k: v for k, v in a.__dict__.items() if not k.startswith("_")}
                        )
                    else:
                        actions.append({"type": getattr(a, "type", "unknown")})

                tool_calls.append(ToolCall(
                    id=item.call_id,
                    name=NATIVE_COMPUTER_USE_TOOL,
                    arguments={"actions": actions},
                ))

            elif item_type == "function_call":
                args = getattr(item, "arguments", "{}")
                if isinstance(args, str):
                    try:
                        args = json.loads(args)
                    except (json.JSONDecodeError, ValueError) as e:
                        logger.warning(
                            f"Failed to parse args for {item.name}: {e}. "
                            f"Raw args (first 300 chars): {repr(args[:300])}"
                        )
                        args = {"__parse_error": str(e), "__raw_preview": args[:500]}
                tool_calls.append(ToolCall(
                    id=item.call_id,
                    name=item.name,
                    arguments=args,
                ))

        content = "\n".join(content_parts) if content_parts else None

        usage: dict[str, int] = {}
        if hasattr(response, "usage") and response.usage:
            input_t = getattr(response.usage, "input_tokens", 0) or 0
            output_t = getattr(response.usage, "output_tokens", 0) or 0
            usage = {
                "prompt_tokens": input_t,
                "completion_tokens": output_t,
                "total_tokens": input_t + output_t,
            }

        finish_reason = "stop"
        if tool_calls:
            finish_reason = "tool_calls"
        if getattr(response, "status", None) == "incomplete":
            finish_reason = "length"

        return LLMResponse(
            content=content,
            tool_calls=tool_calls,
            finish_reason=finish_reason,
            usage=usage,
        )

    def _parse_response(self, response: Any) -> LLMResponse:
        choice = response.choices[0]
        message = choice.message

        tool_calls = []
        if message.tool_calls:
            for tc in message.tool_calls:
                # Skip server-side tool calls (srvtoolu_*).  These are
                # Anthropic server-executed tools (web_search, web_fetch)
                # whose results are already baked into the response text.
                if tc.id and tc.id.startswith(_ANTHROPIC_SERVER_TOOL_ID_PREFIX):
                    logger.debug(
                        f"Skipping server-side tool call: {tc.function.name} ({tc.id})"
                    )
                    continue

                args = tc.function.arguments
                if isinstance(args, str):
                    try:
                        args = json.loads(args)
                    except (json.JSONDecodeError, ValueError) as e:
                        # Tool call arguments were truncated (likely finish_reason=length)
                        # or malformed. Log and pass empty args — the tool will report
                        # a clear error about missing parameters.
                        logger.warning(
                            f"Failed to parse tool call args for {tc.function.name}: {e}. "
                            f"Raw args (first 300 chars): {repr(args[:300])}"
                        )
                        args = {"__parse_error": str(e), "__raw_preview": args[:500]}
                elif args is None:
                    args = {}
                tool_calls.append(
                    ToolCall(id=tc.id, name=tc.function.name, arguments=args)
                )

        usage = {}
        if response.usage:
            usage = {
                "prompt_tokens": response.usage.prompt_tokens or 0,
                "completion_tokens": response.usage.completion_tokens or 0,
                "total_tokens": response.usage.total_tokens or 0,
            }

        return LLMResponse(
            content=message.content,
            tool_calls=tool_calls,
            finish_reason=choice.finish_reason,
            usage=usage,
        )
